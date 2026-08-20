from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Any

from .vault import VaultDenied

OFFICE_KINDS = {
    "docx": "docx",
    "doc": "docx",
    "word": "docx",
    "xlsx": "xlsx",
    "xls": "xlsx",
    "csv": "xlsx",
    "sheet": "xlsx",
    "spreadsheet": "xlsx",
    "excel": "xlsx",
    "pptx": "pptx",
    "ppt": "pptx",
    "slides": "pptx",
    "powerpoint": "pptx",
}


class OfficeError(ValueError):
    pass


def infer_kind(path: str, kind: str | None = None) -> str:
    raw = (kind or Path(path).suffix.lstrip(".") or "").strip().lower()
    mapped = OFFICE_KINDS.get(raw)
    if not mapped:
        raise OfficeError("kind must be docx, xlsx, or pptx")
    return mapped


def ensure_suffix(path: Path, kind: str) -> Path:
    if path.suffix.lower() != f".{kind}":
        return path.with_suffix(f".{kind}")
    return path


def write_office(path: Path, kind: str, *, title: str = "", body: str = "", rows: Any = None, slides: Any = None) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    if kind == "docx":
        _write_docx(path, title, body)
    elif kind == "xlsx":
        _write_xlsx(path, title, body, rows)
    elif kind == "pptx":
        _write_pptx(path, title, body, slides)
    else:
        raise OfficeError(f"unsupported office kind: {kind}")
    return path


def _write_docx(path: Path, title: str, body: str) -> None:
    from docx import Document

    doc = Document()
    if title.strip():
        doc.add_heading(title.strip(), 0)
    for chunk in (body or "").replace("\r\n", "\n").split("\n\n"):
        chunk = chunk.strip("\n")
        if not chunk.strip():
            continue
        first = chunk.split("\n", 1)[0].strip()
        if first.startswith("# "):
            doc.add_heading(first[2:].strip(), 1)
            rest = chunk.split("\n", 1)[1] if "\n" in chunk else ""
            if rest.strip():
                doc.add_paragraph(rest.strip())
            continue
        lines = [line.rstrip() for line in chunk.split("\n")]
        if all(line.strip().startswith(("- ", "* ")) for line in lines if line.strip()):
            for line in lines:
                text = line.strip()[2:].strip()
                if text:
                    doc.add_paragraph(text, style="List Bullet")
            continue
        doc.add_paragraph("\n".join(lines))
    if not title.strip() and not (body or "").strip():
        doc.add_paragraph("")
    doc.save(path)


def _parse_rows(body: str, rows: Any) -> list[list[str]]:
    if isinstance(rows, list) and rows:
        out = []
        for row in rows:
            if isinstance(row, list):
                out.append(["" if cell is None else str(cell) for cell in row])
            else:
                out.append([str(row)])
        return out
    if isinstance(rows, str) and rows.strip():
        body = rows
    text = (body or "").strip()
    if not text:
        return [[]]
    if "|" in text and "\n" in text:
        table = []
        for line in text.splitlines():
            raw = line.strip()
            if not raw or set(raw.replace("|", "").replace("-", "").replace(":", "").strip()) == set():
                continue
            cells = [cell.strip() for cell in raw.strip("|").split("|")]
            table.append(cells)
        if table:
            return table
    reader = csv.reader(io.StringIO(text))
    return [list(row) for row in reader] or [[]]


def _write_xlsx(path: Path, title: str, body: str, rows: Any) -> None:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = (title.strip() or "Sheet1")[:31] or "Sheet1"
    for r_index, row in enumerate(_parse_rows(body, rows), start=1):
        for c_index, cell in enumerate(row, start=1):
            sheet.cell(r_index, c_index, cell)
    book.save(path)


def _parse_slides(title: str, body: str, slides: Any) -> list[tuple[str, str]]:
    if isinstance(slides, list) and slides:
        out = []
        for item in slides:
            if isinstance(item, dict):
                out.append((str(item.get("title") or ""), str(item.get("body") or item.get("text") or "")))
            else:
                out.append(("", str(item)))
        return out or [("", "")]
    chunks = [part.strip() for part in (body or "").replace("\r\n", "\n").split("\n---\n")]
    out = []
    for chunk in chunks:
        if not chunk:
            continue
        lines = chunk.split("\n", 1)
        head = lines[0].lstrip("# ").strip()
        rest = lines[1].strip() if len(lines) > 1 else ""
        out.append((head, rest))
    if title.strip() and not out:
        return [(title.strip(), (body or "").strip())]
    if title.strip() and out and not out[0][0]:
        out[0] = (title.strip(), out[0][1])
    return out or [(title.strip() or "Slide", "")]


def _write_pptx(path: Path, title: str, body: str, slides: Any) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pres = Presentation()
    layout = pres.slide_layouts[1]
    for head, rest in _parse_slides(title, body, slides):
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = head or " "
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = rest
        else:
            box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            frame = box.text_frame
            frame.text = rest
            for paragraph in frame.paragraphs:
                paragraph.font.size = Pt(20)
    if not pres.slides:
        pres.slides.add_slide(layout)
    pres.save(path)


def safe_office_path(path: Path) -> Path:
    if path.exists() and path.is_dir():
        raise VaultDenied("office path is a folder")
    return path
